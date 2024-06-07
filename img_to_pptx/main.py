import os
import json
import subprocess
from contextlib import contextmanager
from pptx import Presentation
from pptx.util import Inches

@contextmanager
def change_dir(destination):
    original_dir = os.getcwd()
    os.chdir(destination)
    try:
        yield
    finally:
        os.chdir(original_dir)



# Function to extract notes from JSON data
def extract_notes(json_data):
    notes = []
    for page in json_data.get('pages', []):
        if 'note' in page:
            notes.append(page['note'])
        else:
            notes.append('')
    return notes


def merge_pngs_to_pptx(folder_path, output_pptx, notes):
    # Create a presentation object
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)

    # Get a list of .png files in the specified folder
    png_files = [f for f in os.listdir(folder_path) if f.endswith('.png')]
    
    # Sort the files by their names
    png_files.sort()

    for (png_file, note) in zip(png_files, notes):
        # Add a blank slide layout
        slide_layout = prs.slide_layouts[6]  # 5 is the index for a blank slide
        slide = prs.slides.add_slide(slide_layout)

        # Calculate the position and size to fill the slide
        left = Inches(0)
        top = Inches(0)
        width = prs.slide_width
        height = prs.slide_height

        # Define the png file path
        png_path = os.path.join(folder_path, png_file)
        print(png_path)
        # Add the png to the slide
        pic = slide.shapes.add_picture(png_path, left, top, width, height)
        # Scale the picture to fit the slide
        pic.width = width
        pic.height = height

        notes_slide = slide.notes_slide
        text_frame = notes_slide.notes_text_frame
        text_frame.text = note


    # Save the presentation to the specified output file
    prs.save(output_pptx)


typst_presentation = '/home/your/presentation/'

output_pptx = 'out.pptx'

with change_dir(typst_presentation):
    subprocess.run(['typst', 'compile', '--format=png', '--ppi=300', 'main.typ', 'out/main-{{n}}.png'])
    subprocess.run(['polylux2pdfpc', 'main.typ', '--font-path=helvetica'])

# Read JSON data from a file
with open(typst_presentation + 'main.pdfpc', 'r') as file:
    data = json.load(file)

# Extract notes
notes = extract_notes(data)

merge_pngs_to_pptx(typst_presentation + 'out/', output_pptx, notes)